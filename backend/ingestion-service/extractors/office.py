"""
Office document extractors.

  extract_docx  — .docx  (python-docx)
  extract_xlsx  — .xlsx  (openpyxl)
  extract_pptx  — .pptx  (python-pptx)
"""

import io


# ── DOCX ──────────────────────────────────────────────────────────────────────

def extract_docx(file_bytes: bytes, filename: str) -> dict:
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(io.BytesIO(file_bytes))
    parts: list[str] = []

    # Main body paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # Mark headings so the chunker can use them as boundaries
            if para.style.name.startswith("Heading"):
                parts.append(f"\n## {text}\n")
            else:
                parts.append(text)

    # Tables
    for table in doc.tables:
        rows: list[str] = []
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            parts.append("\n" + "\n".join(rows) + "\n")

    return {"text": "\n".join(parts), "page_count": 0}


# ── XLSX ──────────────────────────────────────────────────────────────────────

def extract_xlsx(file_bytes: bytes, filename: str) -> dict:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"\n# Sheet: {sheet_name}\n")

        for row in ws.iter_rows(values_only=True):
            cells = [str(v) for v in row if v is not None]
            if cells:
                parts.append("\t".join(cells))

    wb.close()
    return {"text": "\n".join(parts), "page_count": 0}


# ── PPTX ──────────────────────────────────────────────────────────────────────

def extract_pptx(file_bytes: bytes, filename: str) -> dict:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(io.BytesIO(file_bytes))
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = [f"\n--- Slide {slide_num} ---"]

        for shape in slide.shapes:
            # Text frames (title, body, text boxes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        slide_parts.append(line)

            # Tables inside slides
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text_frame.text.strip() for cell in row.cells
                    )
                    if row_text.strip():
                        slide_parts.append(row_text)

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_parts.append(f"[Notes] {notes_text}")

        if len(slide_parts) > 1:   # skip empty slides
            parts.extend(slide_parts)

    return {"text": "\n".join(parts), "page_count": len(prs.slides)}
